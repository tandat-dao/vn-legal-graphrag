#!/usr/bin/env python3
"""Dựng tám slide cải tiến sau phản biện thành .pptx sửa được.

Mọi thứ là hình khối và hộp chữ gốc của PowerPoint — không nhúng ảnh — nên
người dùng đổi được màu, chữ, kích thước từng phần tử.

Toạ độ lấy theo lưới 1280x720 điểm ảnh giống bản HTML, quy đổi 96 điểm ảnh
mỗi inch để khớp khổ 13,333 x 7,5 inch (16:9).
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ── bảng màu, giữ đúng bản HTML ────────────────────────────────────────────
NEN      = "F2F4F8"
MUC      = "1E3050"
MUC_NHAT = "5A6B85"
XANH     = "1554A5"
XANH_S   = "1B7FD9"
XANH_LOT = "D8E7F8"
THE      = "E3E6EB"
THE_NHAT = "EBEDF1"
DO       = "C0392B"
DO_LOT   = "FBEAE8"
VANG     = "B45309"
VANG_LOT = "FDF2E3"
KE       = "9AA6B8"
VIEN     = "C9CFDA"
XAM_TH   = "A9B7CB"      # thanh nền chung
TRANG    = "FFFFFF"

CHU = "Calibri"          # có sẵn trong mọi bản Office, dựng dấu tiếng Việt đúng

def C(h): return RGBColor.from_string(h)
def I(px): return Emu(int(px / 96 * 914400))     # điểm ảnh -> EMU

# lề bản HTML: 4,6cqi x 5,2cqi của 1280
LE_X, LE_Y = 66.6, 58.9
RONG = 1280 - 2 * LE_X


# ── khối dựng cơ bản ───────────────────────────────────────────────────────
def nen_slide(sl, mau=NEN):
    hinh = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, I(1280), I(720))
    hinh.fill.solid(); hinh.fill.fore_color.rgb = C(mau)
    hinh.line.fill.background(); hinh.shadow.inherit = False
    return hinh


def hop(sl, x, y, w, h, mau=None, bo=None, vien=None):
    """Hình chữ nhật, bo góc nếu `bo` (bán kính theo điểm ảnh)."""
    if bo:
        hinh = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, I(x), I(y), I(w), I(h))
        hinh.adjustments[0] = min(bo / min(w, h), 0.5)
    else:
        hinh = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, I(x), I(y), I(w), I(h))
    if mau:
        hinh.fill.solid(); hinh.fill.fore_color.rgb = C(mau)
    else:
        hinh.fill.background()
    if vien:
        hinh.line.color.rgb = C(vien); hinh.line.width = Pt(1)
    else:
        hinh.line.fill.background()
    hinh.shadow.inherit = False
    return hinh


def chu(sl, x, y, w, h, doan, canh=PP_ALIGN.LEFT, doc=MSO_ANCHOR.TOP, dan_dong=None):
    """`doan` là list các đoạn; mỗi đoạn là list run: (text, cỡ, đậm, màu[, cách_chữ])."""
    can = sum(max((r[1] for r in runs), default=12) for runs in doan) * (96 / 72) * 1.42
    h = max(h, can)
    tb = sl.shapes.add_textbox(I(x), I(y), I(w), I(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = doc
    for k, runs in enumerate(doan):
        p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
        p.alignment = canh
        if dan_dong:
            p.line_spacing = dan_dong
        for r in runs:
            txt, co, dam, mau = r[0], r[1], r[2], r[3]
            run = p.add_run(); run.text = txt
            run.font.name = CHU; run.font.size = Pt(co)
            run.font.bold = dam; run.font.color.rgb = C(mau)
            if len(r) > 4 and r[4]:
                run.font._rPr.set("spc", str(int(r[4] * 100)))
    return tb


def duong(sl, x1, y1, x2, y2, mau, day=1.6, mui=False, dut=False):
    ln = sl.shapes.add_connector(2, I(x1), I(y1), I(x2), I(y2))
    ln.line.color.rgb = C(mau); ln.line.width = Pt(day)
    if dut:
        from pptx.enum.dml import MSO_LINE_DASH_STYLE
        ln.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    if mui:
        from pptx.oxml.ns import nsmap
        from lxml import etree
        el = ln.line._get_or_add_ln()
        te = etree.SubElement(el, qn("a:tailEnd"))
        te.set("type", "triangle"); te.set("w", "med"); te.set("len", "med")
    return ln


def tieu_de(sl, text, ke=True):
    chu(sl, LE_X, LE_Y - 6, RONG, 60, [[(text, 30, True, MUC)]])
    if ke:
        hop(sl, LE_X, LE_Y + 52, RONG * .52, 2.4, KE)
    return LE_Y + 82        # y bắt đầu phần thân


def nhan(sl, x, y, w, text, mau=MUC_NHAT):
    return chu(sl, x, y, w, 16, [[(text.upper(), 8.5, True, mau, 1.1)]])


def thanh_doi(sl, x, y, w, ten_runs, nen_pc, them_pc, gia_tri, mau_them=XANH,
              cao=13, w_ten=.40, w_so=.12):
    """Một dòng: nhãn — thanh hai tông — con số."""
    wt, ws = w * w_ten, w * w_so
    wb = w - wt - ws - 24
    chu(sl, x, y + cao / 2 - 8, wt - 10, 20, [ten_runs], doc=MSO_ANCHOR.MIDDLE)
    hop(sl, x + wt, y, wb, cao, "DDE3EC", bo=cao / 2)
    if nen_pc:
        hop(sl, x + wt, y, wb * nen_pc, cao, XAM_TH, bo=cao / 2)
    if them_pc:
        hop(sl, x + wt + wb * nen_pc, y, wb * them_pc, cao, mau_them, bo=cao / 2)
    chu(sl, x + wt + wb + 12, y + cao / 2 - 8, ws, 20,
        [[(gia_tri, 12, True, mau_them if them_pc else MUC)]],
        canh=PP_ALIGN.RIGHT, doc=MSO_ANCHOR.MIDDLE)
    return y + cao + 14


def chu_thich_thanh(sl, x, y, muc):
    """muc = [(màu, nhãn), ...]"""
    cx = x
    for mau, ten in muc:
        hop(sl, cx, y + 2, 11, 8, mau, bo=2)
        t = chu(sl, cx + 16, y - 1, 130, 14, [[(ten, 8.5, False, MUC_NHAT)]])
        cx += 20 + len(ten) * 4.6
    return y + 16


prs = Presentation()
prs.slide_width, prs.slide_height = I(1280), I(720)
TRONG = prs.slide_layouts[6]
def moi(mau=NEN):
    sl = prs.slides.add_slide(TRONG); nen_slide(sl, mau); return sl


# ══════════════════════════════════════════════════════════════════════════
# 1 · phân đoạn
# ══════════════════════════════════════════════════════════════════════════
sl = moi(XANH)
chu(sl, 0, 250, 1280, 24, [[("PHẦN BỔ SUNG", 11, True, "9EC8F2", 1.6)]], canh=PP_ALIGN.CENTER)
chu(sl, 0, 292, 1280, 90, [[("Cải tiến sau phản biện", 54, True, TRANG)]], canh=PP_ALIGN.CENTER)
chu(sl, 320, 404, 640, 60,
    [[("Bốn cơ chế mới ở khâu truy hồi, ba nghiên cứu kiểm chứng, và một lĩnh vực "
       "pháp luật chưa từng có trong kho", 15, False, "C5DDF6")]],
    canh=PP_ALIGN.CENTER, dan_dong=1.45)


# ══════════════════════════════════════════════════════════════════════════
# 2 · vấn đề còn lại
# ══════════════════════════════════════════════════════════════════════════
sl = moi(); y0 = tieu_de(sl, "Vấn đề còn lại sau báo cáo")
CW, CG = (RONG - 2 * 20) / 3, 20
CH = 720 - y0 - LE_Y

# — cột 1: đúng Điều, sai Khoản
x = LE_X
hop(sl, x, y0, CW, CH, THE_NHAT, bo=13)
chu(sl, x + 22, y0 + 20, CW - 44, 24, [[("Đúng Điều, sai Khoản", 14, True, MUC)]])
hop(sl, x + 22, y0 + 56, CW - 44, 176, TRANG, bo=8, vien=VIEN)
chu(sl, x + 38, y0 + 74, 160, 22, [[("Điều 138", 13, True, MUC)]])
hop(sl, x + CW - 118, y0 + 72, 80, 20, "DFF0E0", bo=10)
chu(sl, x + CW - 118, y0 + 75, 80, 16, [[("đúng", 9.5, True, "2E7D32")]], canh=PP_ALIGN.CENTER)
for k, (ten, mau_n, mau_c, ghi) in enumerate([
        ("Khoản 1", "EDEFF3", "8892A4", None),
        ("Khoản 2", DO_LOT,   DO,       "hệ chọn"),
        ("Khoản 3", "EDEFF3", "8892A4", None),
        ("Khoản 5", XANH_LOT, XANH,     "đáp án")]):
    yy = y0 + 104 + k * 30
    hop(sl, x + 38, yy, CW - 76, 24, mau_n, bo=5,
        vien=("E0A9A2" if ten == "Khoản 2" else "8FBBEA" if ten == "Khoản 5" else None))
    chu(sl, x + 48, yy + 4, 120, 18, [[(ten, 11, ghi is not None, mau_c)]])
    if ghi:
        chu(sl, x + CW - 160, yy + 5, 110, 16, [[(ghi, 9.5, True, mau_c)]], canh=PP_ALIGN.RIGHT)
chu(sl, x + 22, y0 + 250, 110, 40, [[("0,806", 21, True, MUC_NHAT)], [("đúng Điều", 9, False, MUC_NHAT)]])
chu(sl, x + 132, y0 + 250, 110, 40, [[("0,737", 21, True, DO)], [("đúng Khoản", 9, False, MUC_NHAT)]])
chu(sl, x + 22, y0 + 302, CW - 44, 60,
    [[("Một điều luật có hàng chục khoản cho những trường hợp khác nhau — ", 10, False, MUC_NHAT),
      ("sai khoản là sai câu trả lời", 10, True, MUC), (".", 10, False, MUC_NHAT)]], dan_dong=1.35)

# — cột 2: bộ lọc loại nhầm
x = LE_X + CW + CG
hop(sl, x, y0, CW, CH, THE_NHAT, bo=13)
chu(sl, x + 22, y0 + 20, CW - 44, 24, [[("Bộ lọc loại nhầm", 14, True, MUC)]])
for k, (tt, mo, nen, mau) in enumerate([
        ("Mốc thời gian quá khứ", "lọc hiệu lực lấy nhầm năm → rỗng", DO_LOT, DO),
        ("Thủ tục ngoài sáu thủ tục", "giám hộ · khai tử · cải chính → rỗng", DO_LOT, DO),
        ("Câu không nêu tỉnh", "gán cứng toàn quốc → mất văn bản tỉnh", "EDEFF3", MUC)]):
    yy = y0 + 60 + k * 62
    hop(sl, x + 22, yy, CW - 44, 54, nen, bo=8)
    chu(sl, x + 36, yy + 10, CW - 72, 18, [[(tt, 11, True, mau)]])
    chu(sl, x + 36, yy + 30, CW - 72, 18, [[(mo, 9.5, False, "A8564C" if mau == DO else MUC_NHAT)]])
hop(sl, x + 22, y0 + 250, CW - 44, 36, MUC, bo=8)
chu(sl, x + 22, y0 + 259, CW - 44, 20, [[("hai kiểu hỏng khác nhau", 11, True, TRANG)]], canh=PP_ALIGN.CENTER)
chu(sl, x + 22, y0 + 302, CW - 44, 60,
    [[("Hai nhóm đầu ", 10, False, MUC_NHAT), ("không tìm được văn bản nào", 10, True, MUC),
      ("; nhóm thứ ba có ngữ cảnh nhưng thiếu đúng phần cần.", 10, False, MUC_NHAT)]], dan_dong=1.35)

# — cột 3: không thấy quy định đã đổi
x = LE_X + 2 * (CW + CG)
hop(sl, x, y0, CW, CH, THE_NHAT, bo=13)
chu(sl, x + 22, y0 + 20, CW - 44, 24, [[("Không thấy quy định đã đổi", 14, True, MUC)]])
hop(sl, x + 22, y0 + 60, CW - 44, 30, "EDEFF3", bo=7, vien=VIEN)
chu(sl, x + 22, y0 + 67, CW - 44, 18, [[("(không có cảnh báo)", 10, False, "96A1B4")]], canh=PP_ALIGN.CENTER)
hop(sl, x + 22, y0 + 102, CW - 44, 128, TRANG, bo=8, vien=VIEN)
chu(sl, x + 40, y0 + 120, CW - 80, 20, [[("Hạn mức là 250 m²/cá nhân", 11, True, MUC)]])
for k, w in enumerate([232, 196, 214]):
    hop(sl, x + 40, y0 + 150 + k * 15, w * (CW - 80) / 250, 7, "E4E8EE", bo=3.5)
hop(sl, x + 40, y0 + 200, 150, 18, XANH_LOT, bo=4)
chu(sl, x + 48, y0 + 203, 140, 14, [[("QĐ 69/2024", 9, True, XANH)]])
chu(sl, x + 22, y0 + 302, CW - 44, 60,
    [[("Người đã nộp hồ sơ trước ngày đổi ", 10, False, MUC_NHAT),
      ("không được báo", 10, True, MUC),
      (" rằng quy định áp cho mình có thể khác.", 10, False, MUC_NHAT)]], dan_dong=1.35)


# ══════════════════════════════════════════════════════════════════════════
# 3 · bốn cơ chế mới
# ══════════════════════════════════════════════════════════════════════════
sl = moi(); y0 = tieu_de(sl, "Bốn cơ chế mới")
chu(sl, LE_X, y0 - 6, RONG * .92, 40,
    [[("Cả bốn đều nằm ở ", 11.5, False, MUC_NHAT), ("Giai đoạn 3 — tìm kiếm lai", 11.5, True, MUC),
      (". Bản thể luận, cách dựng đồ thị và khâu sinh ", 11.5, False, MUC_NHAT),
      ("giữ nguyên", 11.5, True, MUC),
      ("; riêng cơ chế thứ nhất bổ sung một loại quan hệ vào đồ thị lúc nạp dữ liệu.",
       11.5, False, MUC_NHAT)]], dan_dong=1.4)
y0 += 46
OW, OG = (RONG - 18) / 2, 18
OH = (720 - y0 - LE_Y - 16) / 2

def o_dau(x, y, w, ten, so, toi=False, mau_so=None):
    chu(sl, x + 22, y + 18, w - 150, 24, [[(ten, 13, True, TRANG if toi else MUC)]])
    chu(sl, x + w - 172, y + 18, 150, 24,
        [[(so, 14, True, mau_so or ("8FC4F5" if toi else XANH))]], canh=PP_ALIGN.RIGHT)

def o_chan(x, y, w, runs):
    chu(sl, x + 22, y + OH - 46, w - 44, 40, [runs], dan_dong=1.32)

# — 1 · lần theo dẫn chiếu
x, y = LE_X, y0
hop(sl, x, y, OW, OH, MUC, bo=13)
o_dau(x, y, OW, "Lần theo dẫn chiếu", "5.366 quan hệ", toi=True)
chu(sl, x + 22, y + 52, OW - 44, 18,
    [[("“…quy định tại khoản 1 và khoản 2 Điều này”", 10, False, "C2CEE0")]], canh=PP_ALIGN.CENTER)
for k, (ten, day) in enumerate([("Khoản 1", False), ("Khoản 2", False), ("Khoản 3", True)]):
    bx = x + 34 + k * 168
    hop(sl, bx, y + 86, 152, 34, "8FC4F5" if day else None, bo=7,
        vien=None if day else "8FC4F5")
    chu(sl, bx, y + 94, 152, 20, [[(ten, 12, True, "12233D" if day else "8FC4F5")]], canh=PP_ALIGN.CENTER)
# Khoản 3 dẫn ngược về Khoản 1 và Khoản 2: đi vòng lên trên hàng ô,
# nếu nối thẳng thì đường xuyên qua chính ô Khoản 2 ở giữa.
_CX3, _TREN = x + 370, y + 76
duong(sl, _CX3, y + 86, _CX3, _TREN, "8FC4F5")
duong(sl, _CX3, _TREN, x + 110, _TREN, "8FC4F5")
for _cx in (x + 110, x + 278):
    duong(sl, _cx, _TREN, _cx, y + 85, "8FC4F5", mui=True)
chu(sl, x + OW - 190, y + 66, 170, 16, [[("khớp câu hỏi", 9, False, "7E93B2")]], canh=PP_ALIGN.RIGHT)
o_chan(x, y, OW, [("Câu trỏ chỉ chứa địa chỉ, không chứa nội dung — tìm theo độ giống từ ngữ ",
                   10, False, "A9BAD2"), ("không thể lần ra", 10, True, TRANG), (".", 10, False, "A9BAD2")])

# — 2 · cross-encoder
x = LE_X + OW + OG
hop(sl, x, y, OW, OH, XANH_LOT, bo=13)
o_dau(x, y, OW, "Cross-encoder xếp lại trong từng văn bản", "+0,047")
for cot, (nhan_c, ds, mx) in enumerate([
        ("TRƯỚC", ["Điều 12", "Điều 7", "Điều 15", "Điều 9 · Khoản 3"], 3),
        ("SAU",   ["Điều 9 · Khoản 3", "Điều 12", "Điều 7", "Điều 15"], 0)]):
    cx = x + 22 + cot * (OW / 2 + 6)
    nhan(sl, cx, y + 52, 120, nhan_c)
    for k, t in enumerate(ds):
        db = k == mx
        hop(sl, cx, y + 70 + k * 22, OW / 2 - 40, 19, XANH if db else "D7DCE4", bo=4)
        chu(sl, cx + 8, y + 73 + k * 22, OW / 2 - 50, 15,
            [[(t, 9.5, db, TRANG if db else MUC_NHAT)]])
chu(sl, x + OW / 2 - 52, y + 122, 110, 16, [[("đáp án lên đầu", 9, False, MUC_NHAT)]], canh=PP_ALIGN.CENTER)
o_chan(x, y, OW, [("Chỉ đổi thứ tự ", 10, False, "2C4A72"), ("bên trong văn bản đã chọn", 10, True, MUC),
                  (" — thứ tự giữa các văn bản giữ nguyên.", 10, False, "2C4A72")])

# — 3 · mở rộng vùng tìm kiếm
x, y = LE_X, y0 + OH + 16
hop(sl, x, y, OW, OH, THE_NHAT, bo=13)
o_dau(x, y, OW, "Mở rộng vùng tìm kiếm", "+0,053")
for hang, (ten, n_xam, n_xanh) in enumerate([("Trước", 50, 0), ("Sau", 50, 50)]):
    yy = y + 58 + hang * 34
    chu(sl, x + 22, yy - 3, 60, 16, [[(ten, 10, hang == 1, MUC if hang else MUC_NHAT)]])
    for i in range(n_xam):
        hop(sl, x + 86 + i * 8.3, yy, 5, 5, "AEB7C6", bo=2.5)
    for i in range(n_xanh):
        hop(sl, x + 86 + i * 8.3, yy + 10, 5, 5, XANH, bo=2.5)
    chu(sl, x + OW - 70, yy - 4, 50, 18,
        [[(str(n_xam + n_xanh), 12, hang == 1, XANH if hang else MUC_NHAT)]], canh=PP_ALIGN.RIGHT)
o_chan(x, y, OW, [("Gấp đôi số ứng viên lấy ra ở bước tìm ngữ nghĩa, và ", 10, False, MUC_NHAT),
                  ("bỏ hệ số chấm theo độ hiếm khái niệm", 10, True, MUC), (".", 10, False, MUC_NHAT)])

# — 4 · phát hiện quy định đã thay đổi
x = LE_X + OW + OG
hop(sl, x, y, OW, OH, THE_NHAT, bo=13)
o_dau(x, y, OW, "Phát hiện quy định đã thay đổi", "tất định", mau_so=VANG)
hop(sl, x + 22, y + 52, OW - 44, 28, VANG_LOT, bo=7, vien="E8C79A")
hop(sl, x + 36, y + 59, 14, 14, VANG, bo=7)
chu(sl, x + 36, y + 60, 14, 12, [[("!", 8, True, TRANG)]], canh=PP_ALIGN.CENTER)
chu(sl, x + 58, y + 60, OW - 100, 16,
    [[("Quy định về nội dung này đã thay đổi", 10, True, "8A4708")]])
chu(sl, x + 22, y + 94, 90, 16, [[("QĐ 18/2016", 9.5, False, MUC_NHAT)]])
hop(sl, x + 116, y + 95, (OW - 160) * .52, 12, "AEB7C6", bo=6)
chu(sl, x + 22, y + 124, 90, 16, [[("QĐ 69/2024", 9.5, True, MUC)]])
hop(sl, x + 116 + (OW - 160) * .52, y + 125, (OW - 160) * .44, 12, XANH, bo=6)
duong(sl, x + 116 + (OW - 160) * .52, y + 86,
      x + 116 + (OW - 160) * .52, y + 146, VANG, day=1.4, dut=True)
chu(sl, x + 116 + (OW - 160) * .52 - 60, y + 148, 120, 16,
    [[("30/09/2024", 9, True, VANG)]], canh=PP_ALIGN.CENTER)
o_chan(x, y, OW, [("Suy từ đồ thị, ", 10, False, MUC_NHAT), ("không do mô hình phát biểu", 10, True, MUC),
                  (". Hệ không tự kết luận người dùng thuộc quy định nào.", 10, False, MUC_NHAT)])


# ══════════════════════════════════════════════════════════════════════════
# 4 · kết quả truy hồi
# ══════════════════════════════════════════════════════════════════════════
sl = moi(); y0 = tieu_de(sl, "Kết quả truy hồi: 0,737 → 0,853")
hop(sl, LE_X, y0 - 6, RONG, 46, VANG_LOT, bo=8)
hop(sl, LE_X, y0 - 6, 5, 46, VANG)
chu(sl, LE_X + 18, y0 + 2, RONG - 36, 34,
    [[("Đây là thang đo khác với Bảng 4.4 của báo cáo. ", 9.5, True, "8A4708"),
      ("Con số dưới đây là tỉ lệ lấy đúng điều khoản chứa đáp án — chỉ đo khâu truy hồi, "
       "tất định, không gọi mô hình sinh. F1 trích dẫn 0,617 của báo cáo đo cả đường, "
       "gồm cả khâu sinh. Hai con số không so thẳng với nhau được.", 9.5, False, "7A4A12")]],
    dan_dong=1.35)
y0 += 56
KW = (RONG - 40) / 2

# — trái: đường bậc thang
nhan(sl, LE_X, y0, 300, "Cộng dồn từng bước · 123 câu")
BX, BY, BW, BH = LE_X + 34, y0 + 40, KW - 50, 168
LO, HI = .72, .87
gx = lambda k: BX + 16 + k * (BW - 40) / 3
gy = lambda v: BY + BH - (v - LO) / (HI - LO) * BH
for v in (.75, .80, .85):
    hop(sl, BX, gy(v), BW, 1, "E7EBF1")
    chu(sl, BX - 34, gy(v) - 7, 30, 14, [[(f"{v:.2f}".replace(".", ","), 8.5, False, "A6B0C0")]],
        canh=PP_ALIGN.RIGHT)
hop(sl, BX, BY + BH, BW, 1.4, "D5DBE4")
diem = [("trước cải tiến", .737), ("+ dẫn chiếu", .753), ("+ cross-encoder", .800), ("+ mở rộng", .853)]
for k in range(3):
    ln = sl.shapes.add_connector(2, I(gx(k)), I(gy(diem[k][1])), I(gx(k + 1)), I(gy(diem[k + 1][1])))
    ln.line.color.rgb = C(XANH); ln.line.width = Pt(2.6)
for k, (ten, v) in enumerate(diem):
    cuoi = k == 3
    r = 9 if cuoi else 6.5
    d = sl.shapes.add_shape(MSO_SHAPE.OVAL, I(gx(k) - r), I(gy(v) - r), I(r * 2), I(r * 2))
    d.fill.solid(); d.fill.fore_color.rgb = C(XANH if cuoi else TRANG)
    d.line.color.rgb = C(XANH); d.line.width = Pt(2.4); d.shadow.inherit = False
    chu(sl, gx(k) - 60, gy(v) - 34, 120, 22,
        [[(f"{v:.3f}".replace(".", ","), 16 if cuoi else 13, True, MUC)]], canh=PP_ALIGN.CENTER)
    chu(sl, gx(k) - 70, BY + BH + 10, 140, 18, [[(ten, 9.5, False, MUC_NHAT)]], canh=PP_ALIGN.CENTER)

cy = y0 + 246
for k, (so, ten, co) in enumerate([("+0,116", "chênh lệch ghép cặp", 17),
                                   ("0,064–0,171", "khoảng tin cậy 95%", 13.5),
                                   ("p = 0,00007", "kiểm định Wilcoxon", 13.5)]):
    cx = LE_X + k * (KW / 3 + 4)
    hop(sl, cx, cy, KW / 3 - 8, 56, XANH_LOT, bo=8)
    chu(sl, cx + 14, cy + 12, KW / 3 - 36, 22, [[(so, co, True, XANH)]])
    chu(sl, cx + 14, cy + 36, KW / 3 - 36, 14, [[(ten, 8.5, False, "2C4A72")]])
hop(sl, LE_X, cy + 64, KW, 34, THE_NHAT, bo=8)
chu(sl, LE_X + 14, cy + 73, KW - 28, 18,
    [[("26 câu tốt lên", 10, True, XANH), (" · ", 10, False, MUC_NHAT),
      ("4 câu kém đi", 10, True, DO), (" · ", 10, False, MUC_NHAT),
      ("93 câu không đổi", 10, True, MUC)]])

# — phải: thanh theo thách thức
x = LE_X + KW + 40
nhan(sl, x, y0, 400, "Mức tăng trong từng nhóm thách thức")
for k, (ten, v, n) in enumerate([("Đa lĩnh vực", .216, 32), ("Đa tầng văn bản", .153, 30),
                                 ("Đa địa phương", .069, 31), ("Đa phiên bản", .019, 30)]):
    yy = y0 + 34 + k * 56
    db = v > .1
    chu(sl, x, yy, 180, 20, [[(ten, 12, True, MUC)]])
    chu(sl, x, yy + 20, 180, 16, [[(f"{n} câu", 9.5, False, MUC_NHAT)]])
    hop(sl, x + 190, yy + 2, max(v / .225 * (KW - 300), 6), 22, XANH if db else "AEB7C6", bo=6)
    chu(sl, x + 200 + v / .225 * (KW - 300), yy + 5, 90, 18,
        [[(f"+{v:.3f}".replace(".", ","), 12, True, XANH if db else MUC_NHAT)]])
hop(sl, x, y0 + 268, KW, 82, MUC, bo=10)
chu(sl, x + 20, y0 + 286, KW - 40, 60,
    [[("74 trên 123 câu (60%)", 10, True, "8FC4F5"),
      (" đã lấy đúng hoàn toàn từ trước — không còn chỗ cải thiện. Trên 49 câu còn lại, "
       "cải thiện được ", 10, False, "DCE7F5"),
      ("26 câu", 10, True, "8FC4F5"), (", mức tăng trung bình ", 10, False, "DCE7F5"),
      ("+0,317", 10, True, "8FC4F5"), (".", 10, False, "DCE7F5")]], dan_dong=1.4)


# ══════════════════════════════════════════════════════════════════════════
# 5 · cơ chế hay ngữ cảnh
# ══════════════════════════════════════════════════════════════════════════
sl = moi(); y0 = tieu_de(sl, "Cơ chế, hay chỉ là thêm ngữ cảnh?")
hop(sl, LE_X, y0 - 6, RONG, 176, THE_NHAT, bo=13)
MUC5 = [("dẫn chiếu", 3.3, .016), ("cross-encoder", 0.0, .047), ("mở rộng", 0.3, .053)]
for panel, (tieu, lay, dinh) in enumerate([
        ("ĐƠN VỊ THÊM VÀO NGỮ CẢNH", lambda m: m[1], lambda v: f"+{v:.1f}".replace(".", ",")),
        ("MỨC TĂNG BAO PHỦ",          lambda m: m[2], lambda v: f"+{v:.3f}".replace(".", ","))]):
    px0 = LE_X + 26 + panel * (RONG / 2 + 4)
    nhan(sl, px0, y0 + 14, 300, tieu)
    lon = max(lay(m) for m in MUC5)
    for k, m in enumerate(MUC5):
        v = lay(m); hgt = v / lon * 88
        db = v == lon
        bx = px0 + 8 + k * 120
        hop(sl, bx, y0 + 140 - hgt, 96, max(hgt, 3), XANH if db else "AEB7C6", bo=5)
        chu(sl, bx - 12, y0 + 140 - hgt - 24, 120, 20,
            [[(dinh(v), 15 if db else 13, True, XANH if db else MUC_NHAT)]], canh=PP_ALIGN.CENTER)
        chu(sl, bx - 12, y0 + 146, 120, 18, [[(m[0], 10, db, MUC if db else MUC_NHAT)]], canh=PP_ALIGN.CENTER)
    hop(sl, px0, y0 + 140, 344, 1.2, "D5DBE4")
hop(sl, LE_X + RONG / 2 - 4, y0 + 8, 1.2, 150, "D5DBE4")

hop(sl, LE_X, y0 + 186, RONG, 46, MUC, bo=10)
chu(sl, LE_X + 22, y0 + 199, RONG - 44, 22,
    [[("Bước đóng góp ", 11.5, False, "DCE7F5"), ("lớn nhất", 11.5, True, "8FC4F5"),
      (" là bước ", 11.5, False, "DCE7F5"), ("không thêm một đơn vị nào", 11.5, True, "8FC4F5"),
      (" — cross-encoder chỉ xếp lại thứ tự các đơn vị đã chọn.", 11.5, False, "DCE7F5")]])

y = y0 + 254
nhan(sl, LE_X, y, KW + 40, "Nhánh đối chứng cho bước thật sự nạp thêm — thêm đúng bằng số đơn vị đó,")
nhan(sl, LE_X, y + 15, KW + 40, "nhưng chọn theo điểm xếp hạng")
yy = y + 42
yy = thanh_doi(sl, LE_X, yy, KW, [("Không thêm", 11, False, MUC)], .761, 0, "0,761", w_ten=.46)
yy = thanh_doi(sl, LE_X, yy, KW, [("Theo điểm xếp hạng", 11, False, MUC),
                                  (" · đối chứng", 11, False, MUC_NHAT)], .761, .014, "0,775", XANH_S,
                                  w_ten=.46)
yy = thanh_doi(sl, LE_X, yy, KW, [("Theo quan hệ dẫn chiếu", 11, True, MUC)], .761, .036, "0,797", w_ten=.46)
chu_thich_thanh(sl, LE_X, yy, [(XAM_TH, "mức nền chung"), (XANH, "phần tăng thêm")])

x = LE_X + KW + 40
hop(sl, x, y + 30, KW, 116, XANH_LOT, bo=10)
chu(sl, x + 22, y + 50, KW - 44, 80,
    [[("Cùng lượng ngữ cảnh, chọn theo dẫn chiếu hơn chọn theo điểm xếp hạng ", 11.5, False, "2C4A72"),
      ("+0,023", 11.5, True, XANH), (".", 11.5, False, "2C4A72")],
     [("", 6, False, "2C4A72")],
     [("Phần chênh này là đóng góp của ", 11.5, False, "2C4A72"), ("cơ chế", 11.5, True, MUC),
      (", không phải của lượng chữ.", 11.5, False, "2C4A72")]], dan_dong=1.4)


# ══════════════════════════════════════════════════════════════════════════
# 6 · tổng quát hóa
# ══════════════════════════════════════════════════════════════════════════
sl = moi(); y0 = tieu_de(sl, "Tổng quát hóa sang lĩnh vực mới")
chu(sl, LE_X, y0 - 6, RONG * .9, 40,
    [[("Thêm pháp luật ", 11.5, False, MUC_NHAT), ("lao động", 11.5, True, MUC),
      (" — lĩnh vực chưa từng được xây dựng cho. Bộ trích xuất viết cho đất đai chạy trên "
       "văn bản lao động mà không phải sửa dòng nào.", 11.5, False, MUC_NHAT)]], dan_dong=1.4)
y0 += 48
SW = (RONG - 3 * 16) / 4
for k, (so, ten, co) in enumerate([("32 → 36", "văn bản", 22), ("3 → 4", "lĩnh vực", 22),
                                   ("4.549 → 7.208", "điều khoản", 18), ("1.447", "quan hệ dẫn chiếu tạo tự động", 22)]):
    sx = LE_X + k * (SW + 16)
    hop(sl, sx, y0, SW, 74, THE, bo=10)
    chu(sl, sx + 18, y0 + 16, SW - 36, 28, [[(so, co, True, MUC)]])
    chu(sl, sx + 18, y0 + 48, SW - 36, 16, [[(ten, 8.5, False, MUC_NHAT)]])
y0 += 96
KH = 720 - y0 - LE_Y
hop(sl, LE_X, y0, KW, KH, THE_NHAT, bo=13)
nhan(sl, LE_X, y0 + 30, KW, "Có làm hỏng ba lĩnh vực cũ không?", MUC_NHAT)
chu(sl, LE_X, y0 + 52, KW, 90, [[("0", 66, True, XANH)]], canh=PP_ALIGN.CENTER)
chu(sl, LE_X, y0 + 132, KW, 24, [[("câu lệch", 14, True, MUC)]], canh=PP_ALIGN.CENTER)
chu(sl, LE_X + 60, y0 + 164, KW - 120, 48,
    [[("Chạy lại đúng ", 10, False, MUC_NHAT), ("123 câu", 10, True, MUC),
      (" trên kho 36 văn bản, so từng câu với kho 32, ở ", 10, False, MUC_NHAT),
      ("cả bốn cấu hình", 10, True, MUC), (" truy hồi.", 10, False, MUC_NHAT)]],
    canh=PP_ALIGN.CENTER, dan_dong=1.35)

x = LE_X + KW + 40
nhan(sl, x, y0 + 4, KW, "Hệ đạt tới đâu trên lĩnh vực chưa từng thấy?")
yy = y0 + 34
yy = thanh_doi(sl, x, yy, KW, [("Ba lĩnh vực gốc", 11, True, MUC),
                               (" · 123 câu", 11, False, MUC_NHAT)], .737, .116, "0,853")
yy = thanh_doi(sl, x, yy, KW, [("Lao động", 11, True, MUC),
                               (" · 10 câu", 11, False, MUC_NHAT)], .700, .150, "0,850")
yy = chu_thich_thanh(sl, x, yy, [(XAM_TH, "trước cải tiến"), (XANH, "phần tăng thêm")])
hop(sl, x, yy + 10, KW, 66, THE_NHAT, bo=8)
hop(sl, x, yy + 10, 5, 66, XANH)
chu(sl, x + 18, yy + 22, KW - 36, 48,
    [[("Câu ", 10, False, MUC_NHAT), ("LD02", 10, True, MUC),
      (" — điều khoản trả lời xác định phạm vi bằng cách dẫn sang Điều 34 — tăng ", 10, False, MUC_NHAT),
      ("0,50 → 1,00", 10, True, MUC),
      (". Cơ chế dẫn chiếu hoạt động đúng thiết kế trên văn bản chưa từng thấy.", 10, False, MUC_NHAT)]],
    dan_dong=1.3)
chu(sl, x, yy + 88, KW, 40,
    [[("Giới hạn: ", 9, True, MUC),
      ("10 câu là mẫu nhỏ so với 123 câu, nên 0,850 đọc như một chỉ báo. "
       "Tập chuẩn do nhóm soạn, cả hai thành viên đã rà toàn bộ.", 9, False, MUC_NHAT)]], dan_dong=1.3)


# ══════════════════════════════════════════════════════════════════════════
# 7 · kiểm chứng độc lập
# ══════════════════════════════════════════════════════════════════════════
sl = moi(); y0 = tieu_de(sl, "Kiểm chứng độc lập")

nhan(sl, LE_X, y0, KW, "Chấm lại 295 trích dẫn bằng hai bộ chấm khác")
for k, (ten, qh, v, sinh) in enumerate([
        ("Gemini 2.5 Pro", "chính là mô hình sinh", 88.1, True),
        ("Qwen3-4B-Instruct", "khác nhà, khác kiến trúc", 83.7, False),
        ("Gemini 2.5 Flash", "cùng nhà, khác mô hình", 79.0, False)]):
    yy = y0 + 34 + k * 62
    chu(sl, LE_X, yy, 200, 20, [[(ten, 12.5, True, MUC)]])
    chu(sl, LE_X, yy + 21, 200, 16, [[(qh, 9.5, False, VANG if sinh else MUC_NHAT)]])
    bw = KW - 260
    hop(sl, LE_X + 200, yy, bw, 24, "E4E8EE", bo=7)
    hop(sl, LE_X + 200, yy, bw * (v - 70) / 22, 24, VANG if sinh else XANH, bo=7)
    chu(sl, LE_X + 206 + bw, yy + 3, 60, 20,
        [[(f"{v:.1f}".replace(".", ",") + "%", 14, True, VANG if sinh else MUC)]])
chu(sl, LE_X + 200, y0 + 226, KW - 200, 16,
    [[("trục bắt đầu từ 70% để thấy rõ mức chênh", 9, False, MUC_NHAT)]])
hop(sl, LE_X, y0 + 268, KW, 96, VANG_LOT, bo=10)
hop(sl, LE_X, y0 + 268, 5, 96, VANG)
chu(sl, LE_X + 20, y0 + 286, KW - 40, 76,
    [[("Bộ chấm ", 10.5, False, "7A4A12"), ("chính là mô hình sinh", 10.5, True, VANG),
      (" cho điểm cao nhất; hai bộ chấm độc lập đều thấp hơn. Số liệu ", 10.5, False, "7A4A12"),
      ("xác nhận", 10.5, True, MUC),
      (" cảnh báo mà báo cáo đã tự nêu ở mục 4.4.1 — nên trình bày ", 10.5, False, "7A4A12"),
      ("79,0%", 10.5, True, VANG), (" làm cận dưới thận trọng.", 10.5, False, "7A4A12")]], dan_dong=1.4)

x = LE_X + KW + 40
nhan(sl, x, y0, KW, "Hàng thứ năm của Bảng 4.13 · mô hình chạy tại chỗ")
nhan(sl, x + 200, y0 + 30, 120, "NAIVE RAG")
nhan(sl, x + 330, y0 + 30, 120, "GRAPHRAG", XANH)
for k, (ten, nv, gr) in enumerate([("Gemini 2.5 Pro", .435, .617), ("Cục bộ 30B", .317, .583)]):
    yy = y0 + 54 + k * 76
    chu(sl, x, yy + 20, 190, 20, [[(ten, 12.5, True, MUC)]])
    bw = KW - 260
    hop(sl, x + 200, yy, bw * nv, 22, "AEB7C6", bo=6)
    chu(sl, x + 208 + bw * nv, yy + 3, 70, 18, [[(f"{nv:.3f}".replace(".", ","), 11, False, MUC_NHAT)]])
    hop(sl, x + 200, yy + 28, bw * gr, 22, XANH, bo=6)
    chu(sl, x + 208 + bw * gr, yy + 31, 70, 18, [[(f"{gr:.3f}".replace(".", ","), 11, True, XANH)]])
for k, (so, ten) in enumerate([("+0,266", "chênh lệch của mô hình 30B"), ("−0,034", "kém Gemini 2.5 Pro")]):
    cx = x + k * (KW / 2 + 6)
    hop(sl, cx, y0 + 214, KW / 2 - 6, 56, XANH_LOT, bo=8)
    chu(sl, cx + 16, y0 + 226, KW / 2 - 40, 24, [[(so, 17, True, XANH)]])
    chu(sl, cx + 16, y0 + 250, KW / 2 - 40, 14, [[(ten, 8.5, False, "2C4A72")]])
hop(sl, x, y0 + 286, KW, 78, MUC, bo=10)
chu(sl, x + 20, y0 + 304, KW - 40, 60,
    [[("Mô hình chạy ", 10.5, False, "DCE7F5"), ("ngay tại chỗ", 10.5, True, "8FC4F5"),
      (", không gửi dữ liệu ra ngoài. Với hệ xử lý dữ liệu công dân, đây là giá trị thực tế.",
       10.5, False, "DCE7F5")]], dan_dong=1.4)


# ══════════════════════════════════════════════════════════════════════════
# 8 · nút thắt đã dịch chỗ
# ══════════════════════════════════════════════════════════════════════════
sl = moi(); y0 = tieu_de(sl, "Nút thắt đã dịch chỗ")

hop(sl, LE_X, y0, KW, 168, VANG_LOT, bo=13)
nhan(sl, LE_X + 24, y0 + 22, KW - 48, "Bộ sinh trích theo lượng ngữ cảnh nhận được", "8A4708")
for hang, (ten, n_day, n_mo, mau, so) in enumerate([
        ("Hệ trích", 2, 1, VANG, "2,71"), ("Đáp án cần", 1, 1, XANH, "1,57")]):
    yy = y0 + 56 + hang * 46
    chu(sl, LE_X + 24, yy + 4, 110, 22, [[(ten, 12.5, True, MUC)]])
    for i in range(n_day + n_mo):
        hop(sl, LE_X + 140 + i * 52, yy, 42, 28, mau, bo=6)
        if i >= n_day:      # phần lẻ: viền thay vì đặc, thay cho độ mờ
            hop(sl, LE_X + 140 + i * 52, yy, 42, 28, TRANG, bo=6, vien=mau)
    chu(sl, LE_X + 310, yy + 1, 90, 26, [[(so, 20, True, mau)]])
hop(sl, LE_X + 140, y0 + 148, 190, 26, "F7DCD8", bo=13)
chu(sl, LE_X + 140, y0 + 154, 190, 18, [[("thừa 1,14 trích dẫn", 10.5, True, VANG)]], canh=PP_ALIGN.CENTER)

x = LE_X + KW + 40
chu(sl, x, y0 + 10, KW, 56,
    [[("Bao phủ tăng ", 16, True, MUC), ("+0,116", 16, True, XANH),
      (" nhưng F1 của khâu sinh ", 16, True, MUC), ("không thay đổi", 16, True, VANG)]], dan_dong=1.3)
chu(sl, x, y0 + 76, KW, 90,
    [[("Hệ trích trung bình ", 11, False, MUC_NHAT), ("2,71", 11, True, MUC),
      (" điều khoản trong khi đáp án chuẩn chỉ cần ", 11, False, MUC_NHAT), ("1,57", 11, True, MUC),
      (". Mỗi trích dẫn thừa kéo F1 xuống, nên ngữ cảnh tốt hơn ", 11, False, MUC_NHAT),
      ("không tự chuyển thành", 11, True, MUC), (" câu trả lời tốt hơn.", 11, False, MUC_NHAT)],
     [("", 6, False, MUC_NHAT)],
     [("Đo trên tập kiểm thử mới, 28 câu.", 9, False, MUC_NHAT)]], dan_dong=1.45)

y = y0 + 196
nhan(sl, LE_X, y, KW, "Còn bao nhiêu dư địa nếu chọn lọc trích dẫn tối ưu")
yy = y + 30
yy = thanh_doi(sl, LE_X, yy, KW, [("F1 khâu sinh hiện tại", 11, False, MUC)], .571, 0, "0,571", cao=17)
yy = thanh_doi(sl, LE_X, yy, KW, [("Nếu chọn lọc tối ưu", 11, True, MUC)], .571, .171, "0,742", cao=17)

hop(sl, x, y + 12, KW, 88, XANH_LOT, bo=10)
chu(sl, x + 22, y + 26, KW - 44, 34, [[("+0,171", 26, True, XANH)]])
chu(sl, x + 22, y + 62, KW - 44, 30,
    [[("dư địa còn lại — ", 10, False, "2C4A72"),
      ("lớn hơn toàn bộ mức cải thiện của đợt này (+0,116)", 10, True, MUC)]], dan_dong=1.3)

hop(sl, LE_X, 720 - LE_Y - 56, RONG, 56, MUC, bo=10)
chu(sl, LE_X + 22, 720 - LE_Y - 38, RONG - 44, 24,
    [[("Khâu hạn chế nhất đã chuyển từ ", 13, True, TRANG), ("truy hồi", 13, True, "8FC4F5"),
      (" sang ", 13, True, TRANG), ("chọn trích dẫn", 13, True, "8FC4F5")]])

prs.save("slide-cai-tien.pptx")
print("đã lưu slide-cai-tien.pptx ·", len(prs.slides.__iter__.__self__._sldIdLst), "slide")
